"""
Document Parser

Parses various document formats for ingestion
"""

from typing import Dict, Any, Optional, BinaryIO
import logging
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)


class ParsedDocument:
    """Parsed document structure"""
    
    def __init__(
        self,
        content: str,
        metadata: Dict[str, Any],
        file_type: str,
        file_path: Optional[str] = None
    ):
        self.content = content
        self.metadata = metadata
        self.file_type = file_type
        self.file_path = file_path


class DocumentParser:
    """
    Parser for various document formats
    """
    
    def __init__(self):
        self.supported_formats = [
            '.txt', '.md', '.pdf', '.docx', '.doc',
            '.xlsx', '.xls', '.pptx', '.ppt',
            '.html', '.htm', '.csv'
        ]
    
    def parse_file(self, file_path: str) -> ParsedDocument:
        """
        Parse file and extract content
        
        Args:
            file_path: Path to file
        
        Returns:
            ParsedDocument
        """
        
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        suffix = path.suffix.lower()
        
        if suffix not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {suffix}")
        
        logger.info(f"Parsing file: {file_path} ({suffix})")
        
        # Parse based on file type
        if suffix in ['.txt', '.md']:
            content = self._parse_text_file(path)
        elif suffix == '.pdf':
            content = self._parse_pdf(path)
        elif suffix in ['.docx', '.doc']:
            content = self._parse_word(path)
        elif suffix in ['.xlsx', '.xls']:
            content = self._parse_excel(path)
        elif suffix in ['.pptx', '.ppt']:
            content = self._parse_powerpoint(path)
        elif suffix in ['.html', '.htm']:
            content = self._parse_html(path)
        elif suffix == '.csv':
            content = self._parse_csv(path)
        else:
            raise ValueError(f"Parser not implemented for: {suffix}")
        
        # Extract basic metadata
        metadata = self._extract_file_metadata(path)
        
        return ParsedDocument(
            content=content,
            metadata=metadata,
            file_type=suffix,
            file_path=str(path)
        )
    
    def parse_stream(
        self,
        stream: BinaryIO,
        file_type: str,
        filename: Optional[str] = None
    ) -> ParsedDocument:
        """
        Parse file from stream
        
        Args:
            stream: File stream
            file_type: File extension (e.g., '.pdf')
            filename: Original filename
        
        Returns:
            ParsedDocument
        """
        
        # Save stream to temporary file and parse
        import tempfile
        
        with tempfile.NamedTemporaryFile(suffix=file_type, delete=False) as tmp:
            tmp.write(stream.read())
            tmp_path = tmp.name
        
        try:
            result = self.parse_file(tmp_path)
            if filename:
                result.metadata['original_filename'] = filename
            return result
        finally:
            # Cleanup temporary file
            Path(tmp_path).unlink(missing_ok=True)
    
    def _parse_text_file(self, path: Path) -> str:
        """Parse plain text or markdown file"""
        
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return content
    
    def _parse_pdf(self, path: Path) -> str:
        """Parse PDF file"""
        
        try:
            from PyPDF2 import PdfReader
            
            reader = PdfReader(path)
            content_parts = []
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    content_parts.append(f"[Page {page_num}]\n{text}")
            
            return "\n\n".join(content_parts)
        
        except Exception as e:
            logger.error(f"Failed to parse PDF: {e}")
            return f"[Error parsing PDF: {e}]"
    
    def _parse_word(self, path: Path) -> str:
        """Parse Word document"""
        
        try:
            from docx import Document
            
            doc = Document(path)
            content_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    table_text.append(row_text)
                
                if table_text:
                    content_parts.append("\n[Table]\n" + "\n".join(table_text))
            
            return "\n\n".join(content_parts)
        
        except Exception as e:
            logger.error(f"Failed to parse Word document: {e}")
            return f"[Error parsing Word document: {e}]"
    
    def _parse_excel(self, path: Path) -> str:
        """Parse Excel spreadsheet"""
        
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(path, data_only=True)
            content_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"[Sheet: {sheet_name}]")
                
                # Extract data
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter empty rows
                    if any(cell is not None for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        rows.append(row_text)
                
                content_parts.append("\n".join(rows[:100]))  # Limit rows
            
            return "\n\n".join(content_parts)
        
        except Exception as e:
            logger.error(f"Failed to parse Excel: {e}")
            return f"[Error parsing Excel: {e}]"
    
    def _parse_powerpoint(self, path: Path) -> str:
        """Parse PowerPoint presentation"""
        
        try:
            from pptx import Presentation
            
            prs = Presentation(path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                content_parts.append("\n".join(slide_text))
            
            return "\n\n".join(content_parts)
        
        except Exception as e:
            logger.error(f"Failed to parse PowerPoint: {e}")
            return f"[Error parsing PowerPoint: {e}]"
    
    def _parse_html(self, path: Path) -> str:
        """Parse HTML file"""
        
        try:
            from bs4 import BeautifulSoup
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        
        except Exception as e:
            logger.error(f"Failed to parse HTML: {e}")
            return f"[Error parsing HTML: {e}]"
    
    def _parse_csv(self, path: Path) -> str:
        """Parse CSV file"""
        
        try:
            import csv
            
            content_parts = []
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                
                for i, row in enumerate(reader):
                    if i >= 1000:  # Limit rows
                        break
                    
                    row_text = " | ".join(row)
                    content_parts.append(row_text)
            
            return "\n".join(content_parts)
        
        except Exception as e:
            logger.error(f"Failed to parse CSV: {e}")
            return f"[Error parsing CSV: {e}]"
    
    def _extract_file_metadata(self, path: Path) -> Dict[str, Any]:
        """Extract basic file metadata"""
        
        stat = path.stat()
        
        return {
            "filename": path.name,
            "file_size": stat.st_size,
            "file_type": path.suffix,
            "created_date": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_date": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "file_path": str(path.absolute())
        }